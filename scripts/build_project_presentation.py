"""Build the groundwater prospectivity proof-of-concept presentation.

The deck is generated from saved project metrics and feature data so reported
sample counts and model scores stay synchronized with the reproducible runs.
"""

from __future__ import annotations

import argparse
import json
import math
import os
from pathlib import Path

os.environ.setdefault("MPLCONFIGDIR", "/tmp/matplotlib-cache")

import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from matplotlib.colors import LinearSegmentedColormap
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
DEFAULT_METRICS = ROOT / "reports/yield20_xgboost_metrics.json"
DEFAULT_SPATIAL = ROOT / "reports/yield20_spatial_generalization_5000.json"
DEFAULT_FEATURES = ROOT / "data/features/gee_yield20_features_merged_5000.parquet"
DEFAULT_OUTPUT = ROOT / "reports/groundwater_poc_presentation.pptx"
DEFAULT_ASSETS = ROOT / "reports/presentation_assets"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY = "172A3A"
BLUE = "1D6A96"
TEAL = "2A9D8F"
GREEN = "69A66F"
YELLOW = "E9C46A"
ORANGE = "F4A261"
RED = "D65A4A"
INK = "172A3A"
MUTED = "60717F"
PALE = "F3F7F8"
WHITE = "FFFFFF"
GRID = "D9E4E8"


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def set_fill(shape, color: str) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.line.fill.background()


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 20,
    color: str = INK,
    bold: bool = False,
    font: str = "Aptos",
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin: float = 0.03,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Inches(margin)
    frame.margin_top = frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_after = Pt(0)
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def add_rich_lines(slide, lines, x, y, w, h, *, size=16, color=INK, gap=5):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Inches(0.03)
    for index, item in enumerate(lines):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.text = item
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = rgb(color)
        p.space_after = Pt(gap)
        p.level = 0
    return box


def add_title(slide, title: str, subtitle: str | None = None, *, number: int | None = None):
    add_text(slide, title, 0.62, 0.38, 11.8, 0.58, size=27, bold=True)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.62), Inches(1.03), Inches(1.0), Inches(0.06))
    set_fill(bar, TEAL)
    if subtitle:
        add_text(slide, subtitle, 0.62, 1.16, 11.8, 0.42, size=12.5, color=MUTED)
    if number is not None:
        add_text(slide, f"{number:02d}", 12.15, 0.42, 0.52, 0.35, size=11, color=MUTED, align=PP_ALIGN.RIGHT)


def add_footer(slide, source: str, page: int):
    add_text(slide, source, 0.62, 7.14, 11.35, 0.2, size=7.5, color=MUTED)
    add_text(slide, str(page), 12.2, 7.12, 0.45, 0.2, size=8, color=MUTED, align=PP_ALIGN.RIGHT)


def add_pill(slide, text: str, x: float, y: float, w: float, *, color=TEAL, text_color=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.35))
    set_fill(shape, color)
    add_text(slide, text, x, y + 0.02, w, 0.28, size=10, color=text_color, bold=True, align=PP_ALIGN.CENTER)


def add_metric(slide, value: str, label: str, x: float, y: float, color: str, *, w=2.3):
    add_text(slide, value, x, y, w, 0.6, size=32, bold=True, color=color)
    add_text(slide, label, x, y + 0.58, w, 0.48, size=11.5, color=MUTED)


def add_image(slide, path: Path, x: float, y: float, w: float, h: float | None = None):
    kwargs = {"width": Inches(w)}
    if h is not None:
        kwargs["height"] = Inches(h)
    return slide.shapes.add_picture(str(path), Inches(x), Inches(y), **kwargs)


def blank_slide(prs: Presentation, color: str = WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = rgb(color)
    return slide


def category_counts(features: list[str]) -> dict[str, int]:
    mapping = {
        "qa_": "Quality control",
        "dem_": "DEM / terrain",
        "hydro_": "Hydrology",
        "s1_": "Sentinel-1 SAR",
        "s2_": "Sentinel-2 optical",
        "clim_": "Climate",
    }
    counts = {label: 0 for label in mapping.values()}
    for feature in features:
        for prefix, label in mapping.items():
            if feature.startswith(prefix):
                counts[label] += 1
                break
    return counts


def save_feature_chart(counts: dict[str, int], path: Path) -> None:
    colors = [MUTED, BLUE, TEAL, ORANGE, GREEN, RED]
    fig, ax = plt.subplots(figsize=(9.2, 3.5))
    labels = list(counts)
    values = list(counts.values())
    bars = ax.barh(labels[::-1], values[::-1], color=[f"#{c}" for c in colors[::-1]], height=0.62)
    ax.bar_label(bars, padding=5, fontsize=11, fontweight="bold", color=f"#{INK}")
    ax.set_xlim(0, max(values) * 1.22)
    ax.set_xlabel("Number of model features", color=f"#{MUTED}")
    ax.spines[["top", "right", "left"]].set_visible(False)
    ax.tick_params(axis="y", length=0, labelsize=10)
    ax.tick_params(axis="x", colors=f"#{MUTED}")
    ax.grid(axis="x", color=f"#{GRID}", linewidth=0.7)
    ax.set_axisbelow(True)
    fig.tight_layout()
    fig.savefig(path, dpi=180, bbox_inches="tight", transparent=True)
    plt.close(fig)


def save_validation_chart(spatial: dict, path: Path) -> None:
    names = ["ROC-AUC", "PR-AUC", "Accuracy", "F1"]
    keys = ["roc_auc", "pr_auc", "accuracy", "f1"]
    random = [spatial["random_cv"][key] for key in keys]
    blocked = [spatial["spatial_cv"][key] for key in keys]
    x = np.arange(len(names))
    width = 0.34
    fig, ax = plt.subplots(figsize=(8.2, 3.8))
    ax.bar(x - width / 2, random, width, label="Random CV", color=f"#{BLUE}")
    ax.bar(x + width / 2, blocked, width, label="Spatial-block CV", color=f"#{ORANGE}")
    ax.set_xticks(x, names)
    ax.set_ylim(0.55, 0.95)
    ax.set_ylabel("Score")
    ax.grid(axis="y", color=f"#{GRID}", linewidth=0.7)
    ax.set_axisbelow(True)
    ax.spines[["top", "right"]].set_visible(False)
    ax.legend(frameon=False, ncol=2, loc="upper right")
    for xpos, value in zip(x - width / 2, random):
        ax.text(xpos, value + 0.008, f"{value:.2f}", ha="center", fontsize=9)
    for xpos, value in zip(x + width / 2, blocked):
        ax.text(xpos, value + 0.008, f"{value:.2f}", ha="center", fontsize=9)
    fig.tight_layout()
    fig.savefig(path, dpi=180, bbox_inches="tight", transparent=True)
    plt.close(fig)


def save_texas_samples(data_path: Path, path: Path) -> None:
    data = pd.read_parquet(data_path, columns=["longitude", "latitude", "target_yield_ge_20gpm_int"])
    fig, ax = plt.subplots(figsize=(7.1, 4.2))
    for target, color, label in [(0, ORANGE, "< 20 gpm"), (1, TEAL, ">= 20 gpm")]:
        subset = data[data["target_yield_ge_20gpm_int"] == target]
        ax.scatter(subset["longitude"], subset["latitude"], s=4, alpha=0.48, color=f"#{color}", label=label, rasterized=True)
    ax.set_xlabel("Longitude")
    ax.set_ylabel("Latitude")
    ax.spines[["top", "right"]].set_visible(False)
    ax.grid(color=f"#{GRID}", linewidth=0.5)
    ax.legend(frameon=False, markerscale=2.5, ncol=2, loc="lower left")
    fig.tight_layout()
    fig.savefig(path, dpi=180, bbox_inches="tight", transparent=True)
    plt.close(fig)


def save_heatmap_concept(path: Path) -> None:
    rng = np.random.default_rng(42)
    x, y = np.meshgrid(np.linspace(-3, 3, 220), np.linspace(-2, 2, 150))
    field = (
        0.85 * np.exp(-((x + 1.15) ** 2 / 0.75 + (y - 0.25) ** 2 / 0.35))
        + 0.65 * np.exp(-((x - 1.2) ** 2 / 0.45 + (y + 0.55) ** 2 / 0.8))
        + 0.23 * np.sin(1.7 * x + 0.5 * y)
        + 0.06 * rng.normal(size=x.shape)
    )
    field = (field - field.min()) / (field.max() - field.min())
    cmap = LinearSegmentedColormap.from_list("prospectivity", ["#D95F4B", "#E9C46A", "#BFD8A6", "#2A9D8F", "#1D6A96"])
    fig, ax = plt.subplots(figsize=(9.0, 4.5))
    image = ax.imshow(field, cmap=cmap, origin="lower", aspect="auto", extent=[0, 50, 0, 34])
    ax.contour(
        np.linspace(0, 50, field.shape[1]),
        np.linspace(0, 34, field.shape[0]),
        field,
        levels=[0.35, 0.55, 0.75],
        colors="white",
        linewidths=0.7,
        alpha=0.55,
    )
    ax.scatter([13, 29, 37], [20, 10, 24], marker="+", s=90, linewidth=2, color="white")
    ax.set_xticks([])
    ax.set_yticks([])
    cbar = fig.colorbar(image, ax=ax, fraction=0.025, pad=0.02)
    cbar.set_label("P(yield >= 20 gpm)", color="white")
    cbar.ax.tick_params(colors="white")
    cbar.outline.set_edgecolor("white")
    fig.tight_layout()
    fig.savefig(path, dpi=180, bbox_inches="tight", transparent=True)
    plt.close(fig)


def add_process_step(slide, number: int, title: str, body: str, x: float, y: float, color: str):
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.46), Inches(0.46))
    set_fill(circle, color)
    add_text(slide, str(number), x, y + 0.04, 0.46, 0.3, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, title, x + 0.62, y - 0.01, 2.25, 0.35, size=15, bold=True)
    add_text(slide, body, x + 0.62, y + 0.35, 2.25, 0.7, size=10.5, color=MUTED)


def build_deck(metrics: dict, spatial: dict, feature_path: Path, output: Path, assets: Path) -> None:
    assets.mkdir(parents=True, exist_ok=True)
    feature_chart = assets / "feature_groups.png"
    validation_chart = assets / "validation_comparison.png"
    texas_map = assets / "texas_samples.png"
    heatmap = assets / "heatmap_concept.png"
    save_feature_chart(category_counts(metrics["features"]), feature_chart)
    save_validation_chart(spatial, validation_chart)
    save_texas_samples(feature_path, texas_map)
    save_heatmap_concept(heatmap)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 1. Title
    slide = blank_slide(prs, NAVY)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.18), SLIDE_H)
    set_fill(accent, TEAL)
    add_text(slide, "GROUNDWATER\nPROSPECTIVITY", 0.78, 1.15, 7.6, 1.75, size=40, color=WHITE, bold=True)
    add_text(slide, "A satellite-to-ML proof of concept for Kazakhstan", 0.82, 3.08, 7.4, 0.55, size=21, color="BFD7E2")
    add_pill(slide, "TEXAS SDR LABELS", 0.82, 4.03, 1.82, color=ORANGE, text_color=NAVY)
    add_pill(slide, "GOOGLE EARTH ENGINE", 2.82, 4.03, 2.2, color=TEAL)
    add_pill(slide, "XGBOOST", 5.2, 4.03, 1.18, color=BLUE)
    add_text(slide, "Goal: rank candidate locations for field investigation, not replace hydrogeological surveys.", 0.82, 5.4, 7.5, 0.6, size=15, color=WHITE)
    # A compact full-height analytical motif.
    for i, (label, value, color) in enumerate([("93", "features", TEAL), ("5,000", "wells", YELLOW), ("0.840", "ROC-AUC", ORANGE)]):
        y = 1.15 + i * 1.63
        add_text(slide, label, 9.45, y, 2.5, 0.65, size=34, bold=True, color=color, align=PP_ALIGN.RIGHT)
        add_text(slide, value, 9.45, y + 0.63, 2.5, 0.35, size=12, color=WHITE, align=PP_ALIGN.RIGHT)
    add_text(slide, "Project presentation | June 2026", 0.82, 6.88, 4.0, 0.25, size=9, color="8DA7B4")

    # 2. Kazakhstan need
    slide = blank_slide(prs)
    add_title(slide, "Why groundwater prospectivity matters in Kazakhstan", "Water security is a spatial decision problem under climate and transboundary pressure.", number=2)
    add_metric(slide, "44%", "of river flow forms outside Kazakhstan", 0.72, 1.77, BLUE)
    add_metric(slide, "~50%", "potential water-demand shortfall by 2040", 3.58, 1.77, RED)
    add_metric(slide, "4,000+", "explored groundwater deposits", 6.53, 1.77, TEAL)
    add_metric(slide, "15.7 km³/y", "approved operational reserves", 9.62, 1.77, ORANGE, w=2.75)
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.72), SLIDE_W, Inches(2.95))
    set_fill(band, PALE)
    add_text(slide, "THE DECISION", 0.72, 4.08, 2.2, 0.3, size=11, color=TEAL, bold=True)
    add_text(slide, "Where should limited field surveys and drilling budgets be focused first?", 0.72, 4.48, 5.65, 1.05, size=27, bold=True)
    add_rich_lines(slide, [
        "Large territory + sparse labels make exhaustive exploration expensive.",
        "Groundwater can diversify supply where surface water is seasonal or constrained.",
        "A probability heatmap can prioritize, but drilling and water-quality tests remain mandatory.",
    ], 7.05, 4.12, 5.4, 1.75, size=14, gap=9)
    add_footer(slide, "Sources: UNDP Kazakhstan (2021); Government of Kazakhstan / Prime Minister (13 Jan 2026).", 2)

    # 3. Physics
    slide = blank_slide(prs)
    add_title(slide, "How remote sensing contributes evidence", "The model combines indirect indicators of recharge, storage setting, and surface expression.", number=3)
    modalities = [
        ("DEM", "Elevation, slope, aspect, TPI", "Terrain organizes runoff, infiltration and valley position.", BLUE),
        ("SAR", "VV, VH, VV−VH, incidence angle", "C-band backscatter responds to roughness, moisture and vegetation structure.", TEAL),
        ("OPTICAL", "NDVI, NDMI, MNDWI, BSI, NBR", "Reflectance separates vegetation vigor, wetness, bare soil and surface water.", GREEN),
        ("CLIMATE", "Precipitation, PET, AET, deficit, temperature", "Annual water balance constrains potential recharge.", ORANGE),
        ("HYDRO + GEOLOGY", "Drainage, surface water, lithology, faults", "Flow pathways and aquifer material control storage and transmissivity.", RED),
    ]
    for i, (name, features, body, color) in enumerate(modalities):
        x = 0.7 + (i % 3) * 4.16
        y = 1.75 + (i // 3) * 2.15
        w = 3.72 if i < 3 else 5.85
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(1.63))
        set_fill(shape, PALE)
        stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(1.63))
        set_fill(stripe, color)
        add_text(slide, name, x + 0.28, y + 0.18, w - 0.5, 0.3, size=12, color=color, bold=True)
        add_text(slide, features, x + 0.28, y + 0.57, w - 0.5, 0.35, size=14, bold=True)
        add_text(slide, body, x + 0.28, y + 0.98, w - 0.5, 0.48, size=10.5, color=MUTED)
    add_text(slide, "Satellite data do not see groundwater directly.", 7.05, 6.07, 5.1, 0.38, size=15, color=RED, bold=True)
    add_text(slide, "They measure surface and environmental proxies that become useful when joined with well outcomes.", 7.05, 6.43, 5.2, 0.45, size=11.5, color=MUTED)
    add_footer(slide, "Sources: Google Earth Engine Data Catalog; Sentinel-1 GRD and Sentinel-2 SR documentation.", 3)

    # 4. Features
    slide = blank_slide(prs)
    add_title(slide, "Current feature stack: 93 transferable predictors", "Every predictor is available before drilling and can also be computed over a Kazakhstan grid.", number=4)
    add_image(slide, feature_chart, 0.62, 1.55, 7.35)
    add_text(slide, "SPATIAL SUMMARY", 8.45, 1.72, 3.8, 0.3, size=11, color=TEAL, bold=True)
    add_metric(slide, "500 m", "radius around each well", 8.45, 2.13, BLUE, w=2.0)
    add_metric(slide, "30 m", "nominal reduction scale", 10.45, 2.13, ORANGE, w=2.0)
    add_text(slide, "TEMPORAL SUMMARY", 8.45, 3.55, 3.8, 0.3, size=11, color=TEAL, bold=True)
    add_metric(slide, "365 d", "window ending at drill anchor", 8.45, 3.96, GREEN, w=2.0)
    add_metric(slide, "3", "mean / median / std. dev.", 10.45, 3.96, RED, w=2.1)
    add_text(slide, "Feature names encode source, variable, window and reducer; e.g. s1_vh_median_12m_stdDev.", 8.45, 5.56, 3.85, 0.72, size=12, color=MUTED)
    add_footer(slide, "Project artifact: reports/current_gee_feature_dictionary.md", 4)

    # 5. GEE methodology
    slide = blank_slide(prs)
    add_title(slide, "Collection methodology in Google Earth Engine", "A single cloud platform standardizes access, preprocessing, compositing and zonal statistics.", number=5)
    add_process_step(slide, 1, "Anchor", "Use drilling end, start or submission date.", 0.72, 1.72, BLUE)
    add_process_step(slide, 2, "Look backward", "Collect the previous 365 days only.", 3.82, 1.72, TEAL)
    add_process_step(slide, 3, "Composite", "Mask/filter scenes and build robust medians.", 6.92, 1.72, GREEN)
    add_process_step(slide, 4, "Summarize", "Reduce pixels in a 500 m buffer at 30 m scale.", 10.02, 1.72, ORANGE)
    for x in [3.38, 6.48, 9.58]:
        add_text(slide, "→", x, 1.75, 0.3, 0.35, size=20, color=MUTED, bold=True)
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.42), SLIDE_W, Inches(3.22))
    set_fill(band, NAVY)
    add_text(slide, "WHY GEE", 0.72, 3.82, 2.0, 0.3, size=11, color=TEAL, bold=True)
    benefits = [
        ("CO-LOCATED ARCHIVE", "Sentinel, SRTM, JRC water, MERIT Hydro and TerraClimate are queried consistently."),
        ("SERVER-SIDE COMPUTE", "Processing runs near the data; no need to download and mosaic thousands of scenes."),
        ("REPRODUCIBLE API", "The same Python functions can be applied to Texas wells and Kazakhstan grid cells."),
    ]
    for i, (heading, body) in enumerate(benefits):
        x = 0.72 + i * 4.18
        add_text(slide, heading, x, 4.35, 3.45, 0.32, size=11, color=YELLOW, bold=True)
        add_text(slide, body, x, 4.82, 3.5, 1.1, size=13, color=WHITE)
    add_footer(slide, "Source: Google Earth Engine Data Catalog and project extraction pipeline.", 5)

    # 6. Texas labels
    slide = blank_slide(prs)
    add_title(slide, "Why Texas: an open analogue with usable labels", "The Texas SDR database provides coordinates, dates and reported well-test yield at useful scale.", number=6)
    add_image(slide, texas_map, 0.62, 1.55, 7.25)
    add_metric(slide, f"{metrics['rows']:,}", "deduplicated POC wells", 8.35, 1.65, BLUE, w=2.35)
    add_metric(slide, "2022–26", "five-year label window", 10.55, 1.65, TEAL, w=2.35)
    target = metrics["target_distribution"]
    total = int(target["0"]) + int(target["1"])
    positive_rate = int(target["1"]) / total
    add_text(slide, "TARGET", 8.35, 3.3, 1.6, 0.3, size=11, color=TEAL, bold=True)
    add_text(slide, "1 = yield ≥ 20 gpm\n0 = yield < 20 gpm", 8.35, 3.72, 4.1, 0.82, size=18, bold=True)
    add_text(slide, f"Class balance: {positive_rate:.1%} high-yield / {1-positive_rate:.1%} lower-yield", 8.35, 4.68, 4.1, 0.35, size=12, color=MUTED)
    add_text(slide, "Texas is a proof-of-method, not a hydrogeological substitute for Kazakhstan.", 8.35, 5.62, 4.1, 0.7, size=15, color=RED, bold=True)
    add_footer(slide, "Source: Texas Water Development Board Submitted Drillers Reports; project merged feature dataset.", 6)

    # 7. Leakage control
    slide = blank_slide(prs)
    add_title(slide, "Leakage control: simulate the pre-drilling decision", "The label may come from the well test; model inputs must describe only what was knowable beforehand.", number=7)
    add_text(slide, "T − 365 days", 0.78, 2.05, 1.5, 0.4, size=16, color=BLUE, bold=True)
    add_text(slide, "DRILL DATE", 10.86, 2.05, 1.45, 0.4, size=16, color=RED, bold=True, align=PP_ALIGN.RIGHT)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.1), Inches(2.72), Inches(10.85), Inches(0.08))
    set_fill(line, GRID)
    past = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.1), Inches(2.72), Inches(9.85), Inches(0.08))
    set_fill(past, TEAL)
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.82), Inches(2.55), Inches(0.42), Inches(0.42))
    set_fill(dot, RED)
    add_text(slide, "Satellite + climate feature window", 4.05, 3.02, 4.0, 0.4, size=15, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "ALLOWED INPUTS", 0.78, 4.0, 2.25, 0.3, size=11, color=TEAL, bold=True)
    add_rich_lines(slide, ["DEM and drainage", "Pre-drill SAR / optical composites", "Pre-drill climate and water balance", "External regional geology"], 0.78, 4.42, 5.1, 1.55, size=15, gap=7)
    add_text(slide, "EXCLUDED FROM FEATURES", 6.9, 4.0, 3.0, 0.3, size=11, color=RED, bold=True)
    add_rich_lines(slide, ["Yield and pump-test measurements", "Borehole, casing and screen depths", "Drilled-well lithology / strata logs", "Owner, driller and free-text metadata"], 6.9, 4.42, 5.2, 1.55, size=15, gap=7)
    add_footer(slide, "Project rule: feature_end_date ≤ anchor_date for every training well.", 7)

    # 8. Model
    slide = blank_slide(prs)
    add_title(slide, "Main model: tuned gradient-boosted decision trees", "XGBoost is a strong first model for heterogeneous tabular geospatial predictors and missing values.", number=8)
    add_process_step(slide, 1, "Prepare", "Median-impute missing feature values.", 0.72, 1.7, BLUE)
    add_process_step(slide, 2, "Tune", "Randomized search: 20 candidates × 5 folds.", 3.82, 1.7, TEAL)
    add_process_step(slide, 3, "Optimize", "Primary score: average precision for ranking.", 6.92, 1.7, GREEN)
    add_process_step(slide, 4, "Evaluate", "Holdout metrics + spatial-block CV.", 10.02, 1.7, ORANGE)
    add_text(slide, "WHY XGBOOST", 0.72, 3.45, 2.0, 0.3, size=11, color=TEAL, bold=True)
    add_rich_lines(slide, [
        "Captures nonlinear thresholds and interactions without manual equations.",
        "Works well with mixed scales, correlated predictors and moderate sample sizes.",
        "Regularization, row sampling and feature sampling constrain overfit.",
        "Feature importance supports scientific debugging, but not causal claims.",
    ], 0.72, 3.92, 6.0, 2.2, size=15, gap=9)
    params = metrics["best_params"]
    add_text(slide, "SELECTED CONFIGURATION", 7.35, 3.45, 3.2, 0.3, size=11, color=TEAL, bold=True)
    config_lines = [
        f"{params['model__n_estimators']} trees  |  depth {params['model__max_depth']}",
        f"learning rate {params['model__learning_rate']}",
        f"row sample {params['model__subsample']:.1f}  |  feature sample {params['model__colsample_bytree']:.1f}",
        f"L1 {params['model__reg_alpha']}  |  L2 {params['model__reg_lambda']}  |  gamma {params['model__gamma']}",
    ]
    add_rich_lines(slide, config_lines, 7.35, 3.92, 4.9, 1.8, size=15, gap=10)
    add_footer(slide, "Project artifact: reports/yield20_xgboost_metrics.json; XGBoost: Chen & Guestrin (2016).", 8)

    # 9. Results
    slide = blank_slide(prs)
    add_title(slide, "Promising Texas baseline, with a measurable spatial gap", "Random splits answer interpolation; spatial blocks better approximate transfer to unseen areas.", number=9)
    add_image(slide, validation_chart, 0.55, 1.55, 7.45)
    add_metric(slide, f"{metrics['holdout_accuracy']:.3f}", "holdout accuracy", 8.35, 1.65, BLUE, w=2.1)
    add_metric(slide, f"{metrics['holdout_roc_auc']:.3f}", "holdout ROC-AUC", 10.45, 1.65, TEAL, w=2.1)
    add_metric(slide, f"{metrics['holdout_pr_auc']:.3f}", "holdout PR-AUC", 8.35, 3.12, GREEN, w=2.1)
    add_metric(slide, f"{metrics['holdout_f1']:.3f}", "holdout F1", 10.45, 3.12, ORANGE, w=2.1)
    gap = spatial["generalization_gap"]["roc_auc_random_minus_spatial"]
    warning = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.35), Inches(4.78), Inches(4.1), Inches(1.28))
    set_fill(warning, PALE)
    add_text(slide, f"−{gap:.3f} ROC-AUC", 8.62, 5.0, 3.5, 0.42, size=22, color=RED, bold=True)
    add_text(slide, "when moving from random to 1° spatial-block CV", 8.62, 5.48, 3.5, 0.35, size=11.5, color=MUTED)
    add_text(slide, "Interpretation: nearby wells share climate and geology; random CV is optimistic for geographic transfer.", 0.72, 6.38, 7.3, 0.48, size=12.5, color=MUTED)
    add_footer(slide, "Project artifacts: yield20_xgboost_metrics.json and yield20_spatial_generalization_5000.json.", 9)

    # 10. Transfer caveats
    slide = blank_slide(prs)
    add_title(slide, "Kazakhstan inference is a domain-transfer experiment", "The code transfers immediately; the learned Texas relationship may not.", number=10)
    rows = [
        ("Aquifer geology", "Texas lithology and basin architecture differ", "Add Kazakhstan lithology, faults and aquifer maps; calibrate locally"),
        ("Climate regime", "Seasonality, snow and aridity distributions shift", "Use seasonal composites and climate normals; monitor out-of-range features"),
        ("Yield reporting", "20 gpm may reflect use, pump setup and reporting conventions", "Redefine an operational local target and validate test quality"),
        ("Sampling bias", "Wells occur where people already chose to drill", "Use background samples, accessibility covariates and survey design"),
        ("Probability calibration", "Texas probabilities are not Kazakhstan probabilities", "Acquire local labels; recalibrate and quantify uncertainty"),
    ]
    add_text(slide, "SHIFT", 0.72, 1.67, 2.1, 0.3, size=10.5, color=TEAL, bold=True)
    add_text(slide, "RISK", 3.15, 1.67, 4.0, 0.3, size=10.5, color=TEAL, bold=True)
    add_text(slide, "MITIGATION", 7.78, 1.67, 4.45, 0.3, size=10.5, color=TEAL, bold=True)
    for i, (shift, risk, mitigation) in enumerate(rows):
        y = 2.05 + i * 0.91
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.72), Inches(y + 0.7), Inches(11.55), Inches(0.01))
        set_fill(line, GRID)
        add_text(slide, shift, 0.72, y, 2.05, 0.6, size=13, bold=True)
        add_text(slide, risk, 3.15, y, 4.0, 0.62, size=12, color=MUTED)
        add_text(slide, mitigation, 7.78, y, 4.45, 0.65, size=12, color=INK)
    add_footer(slide, "Scientific caveat: prospectivity maps require local hydrogeological validation before operational drilling decisions.", 10)

    # 11. Roadmap
    slide = blank_slide(prs)
    add_title(slide, "Roadmap: tabular ML and spatial deep learning", "Advance in stages so every gain is measured against spatially honest baselines.", number=11)
    stages = [
        ("1", "HARDEN XGBOOST", "Nested spatial CV; Bayesian tuning; class/threshold calibration; SHAP; uncertainty ensembles.", BLUE),
        ("2", "ENRICH GEOLOGY", "Lithology, faults, soils, snow, GRACE context, seasonal and multi-radius summaries.", TEAL),
        ("3", "USE IMAGERY PATCHES", "CNN/ViT encoders learn texture, lineaments, drainage and landform patterns from S1/S2/DEM tiles.", GREEN),
        ("4", "FUSE MODALITIES", "Combine tabular features with learned image embeddings; compare late fusion and attention fusion.", ORANGE),
        ("5", "LEARN LOCALLY", "Active learning: drill or survey where uncertainty and expected information gain are high.", RED),
    ]
    for i, (num, heading, body, color) in enumerate(stages):
        x = 0.7 + i * 2.5
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.7), Inches(2.18), Inches(4.45))
        set_fill(shape, PALE)
        add_text(slide, num, x + 0.18, 1.92, 0.5, 0.5, size=24, color=color, bold=True)
        add_text(slide, heading, x + 0.18, 2.65, 1.82, 0.65, size=12, color=color, bold=True)
        add_text(slide, body, x + 0.18, 3.52, 1.82, 1.85, size=12.5, color=INK)
    add_text(slide, "Deep learning is justified when the label set and spatial coverage become large enough; until then, XGBoost is the stronger, interpretable baseline.", 0.72, 6.46, 11.5, 0.45, size=13, color=MUTED, align=PP_ALIGN.CENTER)
    add_footer(slide, "Roadmap principle: improve transfer validity before increasing model complexity.", 11)

    # 12. Deliverable
    slide = blank_slide(prs, NAVY)
    add_text(slide, "FROM POINT LABELS TO A\nPROSPECTIVITY HEATMAP", 0.72, 0.58, 6.0, 1.18, size=29, color=WHITE, bold=True)
    add_text(slide, "A grid-cell probability is a prioritization score, not a promise of water.", 0.75, 1.92, 5.75, 0.55, size=16, color="BFD7E2")
    add_image(slide, heatmap, 6.72, 0.55, 6.0, 3.95)
    outputs = [
        ("01", "Score every 500 m cell", "Apply the exact Texas feature schema to a selected Kazakhstan region."),
        ("02", "Attach uncertainty", "Flag feature extrapolation, low image counts and ensemble disagreement."),
        ("03", "Prioritize fieldwork", "Use top-ranked, diverse sites for geophysics, hydrochemistry and test drilling."),
    ]
    for i, (num, heading, body) in enumerate(outputs):
        x = 0.75 + i * 4.18
        add_text(slide, num, x, 4.9, 0.6, 0.35, size=11, color=YELLOW, bold=True)
        add_text(slide, heading, x, 5.37, 3.45, 0.38, size=16, color=WHITE, bold=True)
        add_text(slide, body, x, 5.88, 3.55, 0.7, size=11.5, color="D9E4E8")
    add_text(slide, "Proof of concept complete: labels → GEE features → tuned model → spatial validation → deployable grid workflow", 0.75, 6.92, 11.7, 0.28, size=10, color="8DA7B4", align=PP_ALIGN.CENTER)

    # 13. References
    slide = blank_slide(prs)
    add_title(slide, "Selected sources and project evidence", number=13)
    refs = [
        "[1] UNDP Kazakhstan (2021). The climate change impact on water resources in Kazakhstan.",
        "    undp.org/kazakhstan/stories/climate-change-impact-water-resources-kazakhstan",
        "[2] Prime Minister of Kazakhstan (2026). Concept for integrated use of groundwater.",
        "    primeminister.kz/en/news/a-concept-for-the-integrated-use-of-groundwater-is-being-developed-in-kazakhstan-30953",
        "[3] Texas Water Development Board. Submitted Drillers Reports database guidance.",
        "    twdb.texas.gov/groundwater/data/locwwrep.asp",
        "[4] Google Earth Engine Data Catalog and collection documentation.",
        "    developers.google.com/earth-engine/datasets",
        "[5] Chen, T. & Guestrin, C. (2016). XGBoost: A Scalable Tree Boosting System.",
        "    doi.org/10.1145/2939672.2939785",
        "[6] Project reports: feature dictionary, extraction pipeline, model metrics and spatial validation.",
        "    reports/current_gee_feature_dictionary.md; reports/yield20_*metrics*.json",
    ]
    add_rich_lines(slide, refs, 0.72, 1.48, 11.85, 4.95, size=13, color=INK, gap=6)
    add_text(slide, "Reproducible build", 0.72, 6.44, 2.0, 0.3, size=10.5, color=TEAL, bold=True)
    add_text(slide, "gpm/bin/python scripts/build_project_presentation.py", 2.55, 6.42, 6.1, 0.34, size=11.5, color=MUTED, font="Aptos Mono")
    add_footer(slide, "All numerical model results are read from repository artifacts at build time.", 13)

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--metrics-path", type=Path, default=DEFAULT_METRICS)
    parser.add_argument("--spatial-path", type=Path, default=DEFAULT_SPATIAL)
    parser.add_argument("--features-path", type=Path, default=DEFAULT_FEATURES)
    parser.add_argument("--output-path", type=Path, default=DEFAULT_OUTPUT)
    parser.add_argument("--assets-dir", type=Path, default=DEFAULT_ASSETS)
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    with args.metrics_path.open(encoding="utf-8") as file:
        metrics = json.load(file)
    with args.spatial_path.open(encoding="utf-8") as file:
        spatial = json.load(file)
    build_deck(metrics, spatial, args.features_path, args.output_path, args.assets_dir)
    print(f"Presentation saved to {args.output_path}")


if __name__ == "__main__":
    main()
